#!/usr/bin/env python3
"""
fileMerger — Universal File Merger
====================================
Recursively scans a directory, extracts text content from supported file
types, optionally removes duplicate content, and combines everything into
a single Markdown file.

Supported formats: .txt, .md, .json, .csv, .xlsx, .pdf, .docx, .pptx, .ipynb

Usage:
    python fileMerger.py [--directory DIR] [--output FILE] [--csv-limit N]
                         [--json-limit N] [--exclude DIR ...] [--no-dedup]
                         [--verbose] [--quiet]
"""

from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path
from typing import Optional

try:
    from tqdm import tqdm
    HAS_TQDM = True
except ImportError:
    HAS_TQDM = False

import nbformat
import pandas as pd
from docx import Document
from pypdf import PdfReader
from pptx import Presentation


# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

SUPPORTED_EXTENSIONS: set[str] = {
    ".txt",
    ".md",
    ".json",
    ".csv",
    ".xlsx",
    ".pdf",
    ".docx",
    ".pptx",
    ".ipynb",
}

# Directories skipped by default during recursive scan
DEFAULT_EXCLUDE_DIRS: set[str] = {
    ".git",
    ".hg",
    ".svn",
    "__pycache__",
    ".mypy_cache",
    ".pytest_cache",
    ".ruff_cache",
    "node_modules",
    ".venv",
    "venv",
    "env",
    ".env",
    "dist",
    "build",
    "target",
    ".idea",
    ".vscode",
}


# ---------------------------------------------------------------------------
# Deduplication
# ---------------------------------------------------------------------------

def _normalize(text: str) -> str:
    """Return a lowercased, whitespace-collapsed version of *text* for comparison."""
    return re.sub(r"\s+", " ", text.lower()).strip()


class Deduplicator:
    """
    Tracks seen content blocks and reports duplicates.

    Deduplication is intentionally opt-in and scoped per run.  Content is
    compared after lowercasing and whitespace normalisation so that trivially
    equivalent blocks (different capitalisation, extra blank lines) are caught.
    """

    def __init__(self) -> None:
        self._seen: set[str] = set()
        self.duplicates_removed: int = 0

    def is_new(self, text: str) -> bool:
        """Return True and record *text* if it has not been seen before."""
        key = _normalize(text)
        if not key:
            return False
        if key in self._seen:
            self.duplicates_removed += 1
            return False
        self._seen.add(key)
        return True


# ---------------------------------------------------------------------------
# Per-format extractors
# ---------------------------------------------------------------------------

def extract_text(path: Path, dedup: Optional[Deduplicator]) -> str:
    text = path.read_text(encoding="utf-8", errors="ignore").strip()
    if not text:
        return ""
    if dedup is not None and not dedup.is_new(text):
        return ""
    return text


def extract_json(path: Path, dedup: Optional[Deduplicator], json_limit: int) -> str:
    raw = path.read_text(encoding="utf-8", errors="ignore")
    try:
        data = json.loads(raw)
        text = json.dumps(data, indent=2, ensure_ascii=False)
    except json.JSONDecodeError as exc:
        return f"*Could not parse JSON: {exc}*"

    if json_limit > 0 and len(text) > json_limit:
        preview = text[:json_limit]
        note = (
            f"\n\n*JSON truncated: showing first {json_limit:,} of "
            f"{len(text):,} characters.*"
        )
        block = f"```json\n{preview}\n```{note}"
    else:
        block = f"```json\n{text}\n```"

    if dedup is not None and not dedup.is_new(text):
        return ""
    return block


def extract_csv(
    path: Path, dedup: Optional[Deduplicator], csv_limit: int
) -> str:
    try:
        df = pd.read_csv(path)
    except Exception as exc:
        return f"*Could not read CSV: {exc}*"

    if df.empty:
        return ""

    if csv_limit > 0 and len(df) > csv_limit:
        table = df.head(csv_limit).to_markdown(index=False)
        table += f"\n\n*Showing first {csv_limit} of {len(df):,} rows.*"
    else:
        table = df.to_markdown(index=False)

    if dedup is not None and not dedup.is_new(table):
        return ""
    return table


def extract_xlsx(
    path: Path, dedup: Optional[Deduplicator], csv_limit: int
) -> str:
    try:
        xl = pd.ExcelFile(path)
    except Exception as exc:
        return f"*Could not open workbook: {exc}*"

    parts: list[str] = []

    for sheet_name in xl.sheet_names:
        try:
            df = xl.parse(sheet_name)
        except Exception as exc:
            parts.append(f"#### Sheet: {sheet_name}\n\n*Error reading sheet: {exc}*")
            continue

        if df.empty:
            continue

        if csv_limit > 0 and len(df) > csv_limit:
            table = df.head(csv_limit).to_markdown(index=False)
            table += (
                f"\n\n*Showing first {csv_limit} of {len(df):,} rows "
                f"in sheet '{sheet_name}'.*"
            )
        else:
            table = df.to_markdown(index=False)

        if dedup is None or dedup.is_new(table):
            parts.append(f"#### Sheet: {sheet_name}\n\n{table}")

    return "\n\n".join(parts)


def extract_pdf(path: Path, dedup: Optional[Deduplicator]) -> str:
    reader = PdfReader(path)
    parts: list[str] = []

    for i, page in enumerate(reader.pages, start=1):
        text = page.extract_text()
        if not text:
            continue
        text = text.strip()
        if not text:
            continue
        # Deduplicate at the full-document level, not per page, to avoid
        # silently dropping repeated headers/footers.
        parts.append(f"<!-- page {i} -->\n{text}")

    combined = "\n\n".join(parts)
    if not combined:
        return ""
    if dedup is not None and not dedup.is_new(combined):
        return ""
    return combined


def extract_docx(path: Path, dedup: Optional[Deduplicator]) -> str:
    doc = Document(path)
    parts: list[str] = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        if dedup is None or dedup.is_new(text):
            parts.append(text)

    return "\n\n".join(parts)


def extract_pptx(path: Path, dedup: Optional[Deduplicator]) -> str:
    prs = Presentation(path)
    slides: list[str] = []

    for i, slide in enumerate(prs.slides, start=1):
        content: list[str] = []

        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = shape.text.strip()
            if not text:
                continue
            if dedup is None or dedup.is_new(text):
                content.append(text)

        if content:
            slides.append(f"### Slide {i}\n\n" + "\n\n".join(content))

    return "\n\n".join(slides)


def extract_notebook(path: Path, dedup: Optional[Deduplicator]) -> str:
    nb = nbformat.read(path, as_version=4)
    parts: list[str] = []

    for cell in nb.cells:
        content = cell.source.strip()
        if not content:
            continue
        if dedup is not None and not dedup.is_new(content):
            continue

        if cell.cell_type == "markdown":
            parts.append(content)
        elif cell.cell_type == "code":
            parts.append(f"```python\n{content}\n```")

    return "\n\n".join(parts)


# ---------------------------------------------------------------------------
# Dispatcher
# ---------------------------------------------------------------------------

def extract_file(
    path: Path,
    dedup: Optional[Deduplicator],
    csv_limit: int,
    json_limit: int,
) -> str:
    suffix = path.suffix.lower()

    try:
        if suffix in (".txt", ".md"):
            return extract_text(path, dedup)
        if suffix == ".json":
            return extract_json(path, dedup, json_limit)
        if suffix == ".csv":
            return extract_csv(path, dedup, csv_limit)
        if suffix == ".xlsx":
            return extract_xlsx(path, dedup, csv_limit)
        if suffix == ".pdf":
            return extract_pdf(path, dedup)
        if suffix == ".docx":
            return extract_docx(path, dedup)
        if suffix == ".pptx":
            return extract_pptx(path, dedup)
        if suffix == ".ipynb":
            return extract_notebook(path, dedup)
    except Exception as exc:
        return f"*Error processing file: {exc}*"

    return ""


# ---------------------------------------------------------------------------
# File discovery
# ---------------------------------------------------------------------------

def discover_files(
    root: Path,
    output_path: Path,
    extra_excludes: set[str],
) -> list[Path]:
    """Return sorted list of supported files under *root*, skipping excluded dirs."""
    exclude_dirs = DEFAULT_EXCLUDE_DIRS | extra_excludes
    found: list[Path] = []

    for path in root.rglob("*"):
        if not path.is_file():
            continue
        if path.resolve() == output_path:
            continue
        # Skip if any path component is an excluded directory name
        if any(part in exclude_dirs for part in path.parts):
            continue
        if path.suffix.lower() in SUPPORTED_EXTENSIONS:
            found.append(path)

    found.sort()
    return found


# ---------------------------------------------------------------------------
# Output writing
# ---------------------------------------------------------------------------

def write_output(
    output_path: Path,
    root: Path,
    files: list[Path],
    toc: list[str],
    sections: list[str],
    processed: int,
    included: int,
    dedup: Optional[Deduplicator],
) -> None:
    with open(output_path, "w", encoding="utf-8") as out:
        out.write("# Merged Document\n\n")

        out.write("## Table of Contents\n\n")
        out.write("\n".join(toc))
        out.write("\n")

        for section in sections:
            out.write(section)

        out.write("\n---\n")
        out.write("## Statistics\n\n")
        out.write(f"- Files discovered: {len(files)}\n")
        out.write(f"- Files processed: {processed}\n")
        out.write(f"- Files included: {included}\n")
        if dedup is not None:
            out.write(f"- Duplicate blocks removed: {dedup.duplicates_removed}\n")
        else:
            out.write("- Deduplication: disabled\n")


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        prog="fileMerger",
        description=(
            "Recursively scan a directory, extract text from supported files, "
            "and combine everything into a single Markdown document."
        ),
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Supported formats:
  .txt  .md  .json  .csv  .xlsx  .pdf  .docx  .pptx  .ipynb

Examples:
  python fileMerger.py
  python fileMerger.py --directory ./docs --output merged.md
  python fileMerger.py --no-dedup --verbose
  python fileMerger.py --exclude tests fixtures --csv-limit 100
        """,
    )

    parser.add_argument(
        "--directory", "-d",
        default=".",
        metavar="DIR",
        help="Directory to scan (default: current directory)",
    )
    parser.add_argument(
        "--output", "-o",
        default="merged.md",
        metavar="FILE",
        help="Output Markdown filename (default: merged.md)",
    )
    parser.add_argument(
        "--csv-limit",
        type=int,
        default=50,
        metavar="N",
        help="Max rows to include from CSV/XLSX files; 0 = unlimited (default: 50)",
    )
    parser.add_argument(
        "--json-limit",
        type=int,
        default=50_000,
        metavar="N",
        help="Max characters to include from JSON files; 0 = unlimited (default: 50000)",
    )
    parser.add_argument(
        "--exclude",
        nargs="+",
        default=[],
        metavar="DIR",
        help="Additional directory names to exclude (e.g. --exclude tests fixtures)",
    )
    parser.add_argument(
        "--no-dedup",
        action="store_true",
        help="Disable duplicate content removal",
    )
    parser.add_argument(
        "--verbose", "-v",
        action="store_true",
        help="Print each file as it is processed",
    )
    parser.add_argument(
        "--quiet", "-q",
        action="store_true",
        help="Suppress all output except errors",
    )

    return parser


def main() -> None:
    parser = build_parser()
    args = parser.parse_args()

    if args.verbose and args.quiet:
        parser.error("--verbose and --quiet are mutually exclusive")

    root = Path(args.directory).resolve()
    if not root.is_dir():
        print(f"Error: '{root}' is not a directory.", file=sys.stderr)
        sys.exit(1)

    output_path = Path(args.output).resolve()
    extra_excludes = set(args.exclude)
    dedup: Optional[Deduplicator] = None if args.no_dedup else Deduplicator()

    # ── Discover files ──────────────────────────────────────────────────────
    if not args.quiet:
        print(f"Scanning: {root}")

    files = discover_files(root, output_path, extra_excludes)

    if not args.quiet:
        print(f"Found {len(files)} supported file(s).\n")

    # ── Process files ───────────────────────────────────────────────────────
    toc: list[str] = []
    sections: list[str] = []
    processed = 0
    included = 0

    iterator = files
    if HAS_TQDM and not args.quiet and not args.verbose:
        iterator = tqdm(files, desc="Processing", unit="file")

    for file in iterator:
        relative = file.relative_to(root)
        anchor = str(relative).replace("\\", "/")

        if args.verbose:
            print(f"  → {anchor}")

        content = extract_file(file, dedup, args.csv_limit, args.json_limit)
        processed += 1

        if not content.strip():
            if args.verbose:
                print(f"     (empty or duplicate — skipped)")
            continue

        included += 1
        toc.append(f"- [{anchor}](#{anchor.replace('/', '').replace('.', '').replace(' ', '-').lower()})")
        sections.append(f"\n---\n\n## {anchor}\n\n{content}\n")

    # ── Write output ────────────────────────────────────────────────────────
    write_output(
        output_path, root, files, toc, sections,
        processed, included, dedup,
    )

    if not args.quiet:
        print(f"\n✓ Output written to: {output_path}")
        print(f"  Files included : {included}/{processed}")
        if dedup is not None:
            print(f"  Duplicates removed: {dedup.duplicates_removed}")


if __name__ == "__main__":
    main()
